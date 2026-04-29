from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


TITLE_RE = re.compile(r"^(#{1,3})\s+(.*)$")
IMAGE_RE = re.compile(r"^!\[[^\]]*\]\(([^)]+)\)\s*$")


def clean_text(text: str) -> str:
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    return text.strip()


def new_title_slide(prs: Presentation, title: str, subtitle: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle and slide.placeholders and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def new_bullets_slide(prs: Presentation, title: str, bullets: list[tuple[int, str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    body = slide.placeholders[1].text_frame
    body.clear()

    first = True
    for level, text in bullets:
        text = clean_text(text)
        if not text:
            continue
        if first:
            p = body.paragraphs[0]
            first = False
        else:
            p = body.add_paragraph()
        p.text = text
        p.level = max(0, min(level, 4))

    if first:
        body.paragraphs[0].text = ""


def new_image_slide(prs: Presentation, title: str, image_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    if slide.shapes.title:
        slide.shapes.title.text = title

    left = Inches(0.7)
    top = Inches(1.2)
    width = Inches(12.0)
    height = Inches(5.6)
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def is_bullet_line(line: str) -> bool:
    return bool(re.match(r"^(\s*)-\s+.*$", line) or re.match(r"^\s*\d+[.)]\s+.*$", line))


def parse_outline(lines: list[str]) -> list[dict]:
    slides: list[dict] = []
    current_title = "Notes"
    current_bullets: list[tuple[int, str]] = []

    def flush_bullets() -> None:
        nonlocal current_bullets
        if current_bullets:
            slides.append({"type": "bullets", "title": current_title, "bullets": current_bullets[:]})
            current_bullets = []

    idx = 0
    total = len(lines)
    while idx < total:
        raw = lines[idx]
        line = raw.rstrip()
        idx += 1
        if not line:
            continue

        if line.startswith("[Jump to draft section]"):
            continue

        heading = TITLE_RE.match(line)
        if heading:
            flush_bullets()
            level = len(heading.group(1))
            title = clean_text(heading.group(2))
            if level == 1:
                slides.append({"type": "section", "title": title})
                current_title = title
            else:
                current_title = title
            continue

        image = IMAGE_RE.match(line)
        if image:
            flush_bullets()
            image_ref = image.group(1).strip().replace("\\", "/")
            slides.append({"type": "image", "title": current_title, "image_ref": image_ref})
            continue

        if re.match(r"^\|.*\|$", line):
            current_bullets.append((0, line))
            continue

        bullet_match = re.match(r"^(\s*)-\s+(.*)$", line)
        if bullet_match:
            indent = len(bullet_match.group(1))
            level = min(indent // 2, 4)
            current_bullets.append((level, bullet_match.group(2)))
            continue

        numbered_match = re.match(r"^\s*\d+[.)]\s+(.*)$", line)
        if numbered_match:
            current_bullets.append((0, numbered_match.group(1)))
            continue

        if re.fullmatch(r"-{5,}", line):
            flush_bullets()
            continue

        # Heuristic: plain text followed by a bullet block is usually a slide title.
        # This turns "Why not listen to me?" into a title with bullets beneath it.
        next_nonempty = ""
        lookahead = idx
        while lookahead < total:
            probe = lines[lookahead].strip()
            if probe:
                next_nonempty = probe
                break
            lookahead += 1

        if next_nonempty and (
            is_bullet_line(next_nonempty)
            or re.match(r"^!\[[^\]]*\]\(([^)]+)\)\s*$", next_nonempty)
            or re.match(r"^\|.*\|$", next_nonempty)
        ):
            flush_bullets()
            current_title = clean_text(line)
            continue

        current_bullets.append((0, line))

    flush_bullets()
    return slides


def build_presentation(md_path: Path, output_path: Path) -> None:
    lines = md_path.read_text(encoding="utf-8").splitlines()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for item in parse_outline(lines):
        if item["type"] == "section":
            new_title_slide(prs, item["title"])
        elif item["type"] == "bullets":
            new_bullets_slide(prs, item["title"], item["bullets"])
        elif item["type"] == "image":
            image_file = (md_path.parent / item["image_ref"]).resolve()
            if image_file.exists():
                new_image_slide(prs, item["title"], image_file)
            else:
                new_bullets_slide(prs, item["title"], [(0, f"Missing image: {item['image_ref']}")])

    prs.save(output_path)


if __name__ == "__main__":
    source = Path("v3_outline.md")
    output = Path("v3_outline.pptx")
    build_presentation(source, output)
    print(f"Created {output}")
